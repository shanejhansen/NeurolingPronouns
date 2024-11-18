from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN

# List of all 30 sentence pairs with uncontracted 'que'
sentence_pairs = [
    ("Thomas paresse puis il tombe.", "Thomas et Marie paressent puis iel tombent."),
    ("Marie regarde avant que elle saute.", "Marie et Charlie regardent avant que iel sautent."),
    ("Thomas chante et il rigole.", "Thomas et Charlie chantent et iel rigolent."),
    ("Marie bâille et elle rigole.", "Marie et Thomas bâillent et iel rigolent."),
    ("Thomas crie et il pleure.", "Thomas et Marie crient et iel pleurent."),
    ("Marie rigole et elle écoute.", "Marie et Charlie rigolent et iel écoutent."),
    ("Thomas regarde pendant que il nettoie.", "Thomas et Marie regardent pendant que iel nettoient."),
    ("Marie pense avant que elle parle.", "Marie et Charlie pensent avant que iel parlent."),
    ("Thomas rigole puis il écoute.", "Thomas et Marie rigolent puis iel écoutent."),
    ("Marie travaille avant que elle marche.", "Marie et Charlie travaillent avant que iel marchent."),
    ("Thomas cuisine et il travaille.", "Charlie cuisine et iel travaille."),
    ("Marie pense pendant que elle marche.", "Charlie pense pendant que iel marche."),
    ("Thomas rêve et il parle.", "Charlie rêve et iel parle."),
    ("Marie regarde avant que elle joue.", "Charlie regarde avant que iel joue."),
    ("Thomas nettoie et il cuisine.", "Charlie nettoie et iel cuisine."),
    ("Marie chante pendant que elle marche.", "Charlie chante pendant que iel marche."),
    ("Thomas écoute pendant que il nettoie.", "Charlie écoute pendant que iel nettoie."),
    ("Marie paresse et elle écoute.", "Charlie paresse et iel écoute."),
    ("Thomas regarde et il pense.", "Charlie regarde et iel pense."),
    ("Marie mange avant que elle parle.", "Charlie mange avant que iel parle."),

]

# The rest of the script remains the same
# Create a new presentation with widescreen (16:9) slide layout
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Define slide layout
blank_slide_layout = prs.slide_layouts[6]

def add_slide_with_sentences(sentence1, sentence2):
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add a text box
    left = Inches(1)
    top = Inches(2.5)  # Adjusted to center vertically
    width = Inches(14)
    height = Inches(4)  # Reduced height to better fit content
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Add first sentence
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = sentence1
    run.font.size = Pt(44)
    
   

# Add slides for each sentence pair
for pair in sentence_pairs:
    add_slide_with_sentences(pair[0], pair[1])

# Save the presentation
prs.save('french_sentences.pptx')